from pptx import *

path = 'C://Users/Dell/Desktop/\u957f\u591c\u6708.pptx'
savepath = 'C://Users/Dell/Desktop/\u957f\u591c\u6708 - experiment.pptx'

prs = Presentation(path)

# <b>: bold
# <u>: 删除线

QContent = '有 <b>#2[p]</b> 的<u>基础概率</u>使敌方全体受到的伤害提高@ <b>#3[p]</b> #，持续 <b>#4[f]</b> 回合，同时对敌方全体造成等同于海瑟音@ <b>#1[p]</b> #攻击力的物理属性伤害。'
QParams = [
    1.4,
    1.0,
    0.2,
    3.0
]

def format_skill_content(content:str, params:list[float|int]) -> str:
    ans = content.replace('<b>', '')
    ans = ans.replace('</b>', '')
    ans = ans.replace('<u>', '')
    ans = ans.replace('</u>', '')
    ans = ans.replace('@', '')
    for i, v in enumerate(params):
        ans = ans.replace(f"#{i + 1}[p]", f"{int(v * 100)}%")
        ans = ans.replace(f"#{i + 1}[f]", f"{int(v)}")
    ans = ans.replace('#', '')
    return ans

def simple_text_frame_replace(text_frame, old_text:str, new_text:str) -> None:
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            run.text = run.text.replace(old_text, new_text)

for slide in prs.slides:
    for shape in slide.shapes:
        if shape.has_text_frame:
            text_frame = shape.text_frame
            text = text_frame.text
            if 'Q About to fill' in text:
                simple_text_frame_replace(text_frame, 'Q About to fill', format_skill_content(QContent, QParams))

prs.save(savepath)